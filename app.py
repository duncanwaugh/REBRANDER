import streamlit as st
from io import BytesIO
from pathlib import Path
from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from openpyxl import load_workbook
import zipfile
from PIL import Image
import imagehash
import base64
import os
import subprocess, tempfile
from docx.opc.constants import RELATIONSHIP_TYPE as RT

def wmf_to_png_blob(wmf_blob: bytes) -> bytes:
    # write temp WMF
    with tempfile.NamedTemporaryFile(suffix=".wmf", delete=False) as wmftmp:
        wmftmp.write(wmf_blob)
    png_path = wmftmp.name + ".png"
    # convert via ImageMagick CLI
    subprocess.run(["convert", wmftmp.name, png_path], check=True)
    # read back PNG
    with open(png_path, "rb") as f:
        return f.read()

# Directory containing old logos (relative to this script)
OLD_LOGO_DIR = Path(__file__).parent / "old_logos"
# Hamming distance threshold for perceptual hash matching
HASH_THRESHOLD = 25

# Load and hash old logos once at startup
def load_old_logo_hashes(threshold: int = HASH_THRESHOLD):
    logo_hashes = []
    if OLD_LOGO_DIR.exists() and OLD_LOGO_DIR.is_dir():
        for img_path in OLD_LOGO_DIR.iterdir():
            if img_path.suffix.lower() in {".png", ".jpg", ".jpeg"}:
                try:
                    img = Image.open(img_path)
                    logo_hashes.append(imagehash.phash(img))
                except Exception:
                    continue
    return logo_hashes

# Text replacement in Word documents
def replace_text_docx(doc: Document, mappings: dict):
    for p in doc.paragraphs:
        for find, replace in mappings.items():
            if find in p.text:
                for run in p.runs:
                    run.text = run.text.replace(find, replace)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for find, replace in mappings.items():
                    if find in cell.text:
                        cell.text = cell.text.replace(find, replace)

def replace_images_docx(doc: Document, old_hashes: list, new_logo_blob: bytes):
    for part in doc.part.package.parts:
        partname = getattr(part, "partname", "").lower()
        ctype    = part.content_type

        if not ctype.startswith("image/"):
            continue

        # 1) grab & convert raw bytes if WMF
        raw = part.blob
        if partname.endswith(".wmf"):
            try:
                raw = wmf_to_png_blob(raw)
                part._content_type = "image/png"
                st.write(f"[DEBUG] Converted WMF '{partname}' → PNG + set content_type")
            except Exception as e:
                st.write(f"[DEBUG] WMF→PNG failed for '{partname}': {e}")
                continue

        # 2) hash it
        try:
            img = Image.open(BytesIO(raw))
        except Exception as e:
            st.write(f"[DEBUG] Cannot open image '{partname}': {e}")
            continue

        h = imagehash.phash(img)
        distances = [abs(h - old_h) for old_h in old_hashes]
        min_dist  = min(distances) if distances else None
        st.write(f"[DEBUG] '{partname}' → min distance = {min_dist}")

        # 3) if it matches, overwrite *both* the part and the rel‐target
        if min_dist is not None and min_dist <= HASH_THRESHOLD:
            st.write(f"[DEBUG] Replacing '{partname}' (distance {min_dist})")
            # a) update the OPC part
            part._blob = new_logo_blob
            part._content_type = "image/png"
            # b) update the relationship target
            for rel in doc.part.rels.values():
                if rel.reltype == RT.IMAGE and rel._target.partname.lower() == partname:
                    rel._target._blob = new_logo_blob
                    rel._target._content_type = "image/png"



# Process .docx files
def process_docx(uploaded_file, mappings: dict, new_logo_bytes: bytes, old_hashes: list):
    doc = Document(uploaded_file)
    replace_text_docx(doc, mappings)
    replace_images_docx(doc, old_hashes, new_logo_bytes)
    output = BytesIO()
    doc.save(output)
    output.seek(0)
    return output

# Process .pptx files with perceptual hash image replacement and debug info (no scaling)
def process_pptx(uploaded_file, mappings: dict, new_logo_bytes: bytes, old_hashes: list):
    prs = Presentation(uploaded_file)
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            # Text replacement
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    for run in p.runs:
                        for find, replace in mappings.items():
                            run.text = run.text.replace(find, replace)
            # Image replacement
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img = Image.open(BytesIO(shape.image.blob))
                    h = imagehash.phash(img)
                    distances = [abs(h - old_h) for old_h in old_hashes]
                    min_dist = min(distances) if distances else None
                    st.write(f"[DEBUG] PPTX slide {slide_idx} image - min Hamming distance: {min_dist}")
                    if min_dist is not None and min_dist <= HASH_THRESHOLD:
                        st.write(f"[DEBUG] Replacing PPTX slide {slide_idx} image (distance {min_dist})")
                        left, top, width, height = shape.left, shape.top, shape.width, shape.height
                        slide.shapes._spTree.remove(shape._element)
                        prs.slides[slide_idx-1].shapes.add_picture(
                            BytesIO(new_logo_bytes), left, top, width, height
                        )
                except Exception as e:
                    st.write(f"[DEBUG] Error processing PPTX slide {slide_idx} image: {e}")
                    continue
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output

# Process .xlsx files: text + image replacement in xl/media with debug info (no scaling)
def process_excel(uploaded_file, mappings: dict, new_logo_bytes: bytes, old_hashes: list):
    wb = load_workbook(filename=BytesIO(uploaded_file.read()))
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=False):
            for cell in row:
                if isinstance(cell.value, str):
                    for find, replace in mappings.items():
                        if find in cell.value:
                            cell.value = cell.value.replace(find, replace)
    interim = BytesIO()
    wb.save(interim)
    interim.seek(0)
    in_zip = zipfile.ZipFile(interim, 'r')
    out_io = BytesIO()
    with zipfile.ZipFile(out_io, 'w') as out_zip:
        for item in in_zip.infolist():
            data = in_zip.read(item.filename)
            if item.filename.startswith('xl/media/') and data:
                try:
                    img = Image.open(BytesIO(data))
                    h = imagehash.phash(img)
                    distances = [abs(h - old_h) for old_h in old_hashes]
                    min_dist = min(distances) if distances else None
                    st.write(f"[DEBUG] Excel media '{item.filename}' - min Hamming distance: {min_dist}")
                    if min_dist is not None and min_dist <= HASH_THRESHOLD:
                        st.write(f"[DEBUG] Replacing Excel media '{item.filename}' (distance {min_dist})")
                        data = new_logo_bytes
                except Exception as e:
                    st.write(f"[DEBUG] Error processing Excel media '{item.filename}': {e}")
            out_zip.writestr(item, data)
    out_io.seek(0)
    return out_io

# --- Streamlit UI Styling (Aecon Lessons Learned style) ---
st.set_page_config(page_title="File Rebrander", page_icon="📘", layout="wide")

# Display logo if available
logo_path = OLD_LOGO_DIR / "logo_1.PNG"
if logo_path.exists():
    st.image(str(logo_path), width=300)
else:
    st.write("Logo file not found; please add 'old_logos/aecon_logo.png'.")

st.markdown("""
<style>
  .stApp { background:#fff; }
  h1,h2 { color:#c8102e; }
  .stButton>button, .stDownloadButton>button { background:#c8102e; color:#fff; }
  body { font-family:'Segoe UI', sans-serif; }
</style>
""", unsafe_allow_html=True)

st.title("File Rebrander")
# Main inputs
mapping_text = st.text_area(
    "Find → Replace mappings (one per line, comma-separated)",
    "Aecon Group Inc. (AGI),North End Connectors (NEC)\nAecon Group Inc.,North End Connectors (NEC)\nAGI,NEC\nAecon,North End Connectors (NEC)",
    height=150
)
mappings = dict(line.split(",",1) for line in mapping_text.splitlines() if line.strip())
new_logo = st.file_uploader("Upload new logo image", type=["png","jpg","jpeg"])
uploaded = st.file_uploader("Upload document(s) to rebrand", type=["docx","pptx","xlsx"],accept_multiple_files=True)

if uploaded and st.button("Rebrand Document(s)"):
    if not new_logo:
        st.error("Please upload the new logo image.")
    else:
        old_hashes    = load_old_logo_hashes()
        new_logo_bytes = new_logo.read()

        # create a ZIP in memory
        zip_buffer = BytesIO()
        with zipfile.ZipFile(zip_buffer, "w") as zf:
            for file in uploaded:
                ext = file.name.split('.')[-1].lower()
                if ext == 'docx':
                    out = process_docx(file, mappings, new_logo_bytes, old_hashes)
                elif ext == 'pptx':
                    out = process_pptx(file, mappings, new_logo_bytes, old_hashes)
                elif ext in ('xlsx','xlsm'):
                    out = process_excel(file, mappings, new_logo_bytes, old_hashes)
                else:
                    continue

                # add to ZIP under a new filename
                zf.writestr(f"rebranded_{file.name}", out.getvalue())

        zip_buffer.seek(0)
        st.success("✅ Batch rebranding complete!")
        st.download_button(
            "📥 Download All Rebranded Files", 
            data=zip_buffer.getvalue(), 
            file_name="rebranded_documents.zip", 
            mime="application/zip"
        )

st.markdown("""
<hr style='border:none;height:2px;background:#c8102e;'/>
<div style='text-align:center;padding:10px;background:#c8102e;color:#fff;'>
  Built for Aecon | Still in development 
</div>
""", unsafe_allow_html=True)
